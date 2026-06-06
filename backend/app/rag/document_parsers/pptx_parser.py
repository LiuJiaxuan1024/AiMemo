from __future__ import annotations

from pathlib import Path

from app.rag.document_parsers.base import DocumentBlock, DocumentImageAsset, DocumentParseError
from app.rag.document_parsers.base import ParsedDocument, normalize_text


def parse_pptx(path: Path) -> ParsedDocument:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise DocumentParseError("DOCUMENT_READER_DEPENDENCY_MISSING", "缺少 python-pptx 依赖，无法解析 PPTX。") from exc

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise DocumentParseError("DOCUMENT_PARSE_FAILED", f"PPTX 解析失败：{exc}") from exc

    blocks: list[DocumentBlock] = []
    image_assets: list[DocumentImageAsset] = []
    title = ""
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_title = _slide_title(slide)
        heading_path = [slide_title] if slide_title else [f"Slide {slide_index}"]
        if slide_title:
            if not title:
                title = slide_title
            blocks.append(
                DocumentBlock(
                    text=slide_title,
                    block_type="heading",
                    heading_path=list(heading_path),
                    page_number=slide_index,
                    metadata={"level": 1, "slide": slide_index},
                )
            )

        shape_texts: list[str] = []
        image_index = 0
        for shape in slide.shapes:
            shape_texts.extend(_shape_texts(shape, MSO_SHAPE_TYPE))
            for image_asset in _shape_image_assets(shape, MSO_SHAPE_TYPE, slide_index, heading_path, image_index):
                image_index += 1
                image_assets.append(image_asset)
        for offset, text in enumerate(_dedupe_preserve_order(shape_texts)):
            if slide_title and text == slide_title:
                continue
            blocks.append(
                DocumentBlock(
                    text=text,
                    block_type="paragraph",
                    heading_path=list(heading_path),
                    page_number=slide_index,
                    source_offset=offset,
                    metadata={"slide": slide_index},
                )
            )

        notes_text = _notes_text(slide)
        if notes_text:
            blocks.append(
                DocumentBlock(
                    text=notes_text,
                    block_type="note",
                    heading_path=[*heading_path, "Notes"],
                    page_number=slide_index,
                    metadata={"slide": slide_index},
                )
            )

    if not blocks:
        raise DocumentParseError("DOCUMENT_TEXT_EMPTY", "未能从 PPTX 中提取到文本。它可能主要由图片组成。")
    return ParsedDocument(parser="pptx", blocks=blocks, title=title or path.stem, image_assets=image_assets)


def _slide_title(slide) -> str:
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is None:
        return ""
    return normalize_text(getattr(title_shape, "text", "") or "")


def _shape_texts(shape, shape_type_enum) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = _text_frame_text(shape.text_frame)
        if text:
            texts.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [normalize_text(cell.text).replace("\n", " ") for cell in row.cells]
            line = " | ".join(cell for cell in cells if cell)
            if line:
                texts.append(line)
    if getattr(shape, "shape_type", None) == shape_type_enum.GROUP:
        for child_shape in shape.shapes:
            texts.extend(_shape_texts(child_shape, shape_type_enum))
    return texts


def _shape_image_assets(shape, shape_type_enum, slide_index: int, heading_path: list[str], image_offset: int) -> list[DocumentImageAsset]:
    assets: list[DocumentImageAsset] = []
    if getattr(shape, "shape_type", None) == shape_type_enum.PICTURE:
        asset_index = image_offset + 1
        alt_text = getattr(shape, "alternative_text", None) or getattr(shape, "name", None)
        image = getattr(shape, "image", None)
        assets.append(
            DocumentImageAsset(
                asset_id=f"pptx-slide-{slide_index}-image-{asset_index}",
                parser="pptx",
                location_label=f"Slide {slide_index} 图片 {asset_index}",
                heading_path=list(heading_path),
                page_number=slide_index,
                source_offset=asset_index,
                alt_text=normalize_text(str(alt_text or "")) or None,
                data=getattr(image, "blob", None) if image is not None else None,
                mime_type=getattr(image, "content_type", None) if image is not None else None,
                width=int(getattr(shape, "width", 0) or 0) or None,
                height=int(getattr(shape, "height", 0) or 0) or None,
            )
        )
    if getattr(shape, "shape_type", None) == shape_type_enum.GROUP:
        nested_offset = image_offset + len(assets)
        for child_shape in shape.shapes:
            child_assets = _shape_image_assets(child_shape, shape_type_enum, slide_index, heading_path, nested_offset)
            nested_offset += len(child_assets)
            assets.extend(child_assets)
    return assets


def _text_frame_text(text_frame) -> str:
    lines: list[str] = []
    for paragraph in text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if text:
            lines.append(text)
    return normalize_text("\n".join(lines))


def _notes_text(slide) -> str:
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return ""
    notes_text_frame = getattr(notes_slide, "notes_text_frame", None)
    if notes_text_frame is None:
        return ""
    return _text_frame_text(notes_text_frame)


def _dedupe_preserve_order(items: list[str]) -> list[str]:
    seen: set[str] = set()
    result: list[str] = []
    for item in items:
        normalized = normalize_text(item)
        if not normalized or normalized in seen:
            continue
        seen.add(normalized)
        result.append(normalized)
    return result
