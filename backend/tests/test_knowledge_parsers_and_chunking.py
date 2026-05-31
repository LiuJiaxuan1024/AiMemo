from pathlib import Path

from docx import Document

from app.rag.document_parsers import parse_document_file
from app.rag.knowledge_chunking import build_chunk_drafts


def test_text_parser_and_chunker_create_basic_drafts(tmp_path: Path) -> None:
    path = tmp_path / "notes.txt"
    path.write_text("第一段内容。\n\n第二段内容。", encoding="utf-8")

    parsed = parse_document_file(path)
    drafts = build_chunk_drafts(parsed.blocks)

    assert parsed.parser == "text"
    assert [block.text for block in parsed.blocks] == ["第一段内容。", "第二段内容。"]
    assert len(drafts) == 1
    assert "第一段内容" in drafts[0].text
    assert drafts[0].token_count > 0
    assert drafts[0].content_hash


def test_markdown_parser_preserves_heading_paths(tmp_path: Path) -> None:
    path = tmp_path / "guide.md"
    path.write_text("# 总览\n\n说明。\n\n## 安装\n\n步骤一。", encoding="utf-8")

    parsed = parse_document_file(path)
    drafts = build_chunk_drafts(parsed.blocks)

    assert parsed.parser == "markdown"
    assert parsed.title == "总览"
    assert any(block.block_type == "heading" for block in parsed.blocks)
    assert any(draft.heading_path == ["总览", "安装"] for draft in drafts)


def test_docx_parser_reads_paragraphs_and_tables(tmp_path: Path) -> None:
    path = tmp_path / "manual.docx"
    document = Document()
    document.add_heading("手册", level=1)
    document.add_paragraph("正文内容。")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "参数"
    table.cell(0, 1).text = "含义"
    document.save(path)

    parsed = parse_document_file(path)

    assert parsed.parser == "docx"
    assert parsed.title == "手册"
    assert any(block.text == "正文内容。" for block in parsed.blocks)
    assert any(block.text == "参数 | 含义" for block in parsed.blocks)


def test_pptx_parser_reads_slides_tables_and_notes(tmp_path: Path) -> None:
    from pptx import Presentation

    path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "路线图"
    slide.placeholders[1].text = "第一阶段：导入资料\n第二阶段：挂载对话"
    table_shape = slide.shapes.add_table(rows=1, cols=2, left=0, top=0, width=1000000, height=300000)
    table_shape.table.cell(0, 0).text = "模块"
    table_shape.table.cell(0, 1).text = "知库"
    slide.notes_slide.notes_text_frame.text = "演讲者备注：强调挂载边界"
    presentation.save(path)

    parsed = parse_document_file(path)

    assert parsed.parser == "pptx"
    assert parsed.title == "路线图"
    assert any(block.block_type == "heading" and block.text == "路线图" for block in parsed.blocks)
    assert any("第一阶段" in block.text for block in parsed.blocks)
    assert any(block.text == "模块 | 知库" for block in parsed.blocks)
    assert any("演讲者备注" in block.text for block in parsed.blocks)
    assert all(block.page_number == 1 for block in parsed.blocks)


def test_pdf_parser_reads_page_text(tmp_path: Path) -> None:
    import fitz

    path = tmp_path / "paper.pdf"
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "PDF 知库内容")
    document.save(path)
    document.close()

    parsed = parse_document_file(path)

    assert parsed.parser == "pdf"
    assert parsed.blocks[0].page_number == 1
    assert "PDF" in parsed.blocks[0].text
